import os
import requests
from pptx import Presentation

API_KEY = os.getenv("SAMBA_API_KEY", "540f8914-997e-46c6-829a-ff76f5d4d265")
API_URL = os.getenv("SAMBA_API_URL", "https://api.sambanova.ai/v1/chat/completions")

def extract_text_from_ppt(ppt_path: str) -> str:
    """Extracts text from a PPT or PPTX file using python-pptx."""
    text = ""
    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        raise Exception(f"Error extracting PPT: {e}")
    return text

def summarize_text(text: str) -> str:
    """Uses SambaNova API to summarize the given text."""
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": "llama3-405b",
        "messages": [{"role": "user", "content": text}],
        "max_tokens": 1024,
        "temperature": 0.7,
        "top_k": 1,
        "top_p": 0.01
    }

    response = requests.post(API_URL, json=payload, headers=headers, timeout=30)

    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    else:
        raise Exception(f"Error: {response.status_code} - {response.text}")

if __name__ == "__main__":
    ppt_path = input("Enter the path to the PPT file: ").strip()

    try:
        ppt_text = extract_text_from_ppt(ppt_path)
        print("PPT Text Extracted Successfully.")
    except Exception as e:
        print(f"Failed to extract PPT text: {e}")
        exit(1)

    try:
        summary = summarize_text(ppt_text)
        print("Summary:\n", summary)
    except Exception as e:
        print(f"Failed to summarize text: {e}")
