import os
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from langchain_community.vectorstores import FAISS
from langchain_core.vectorstores import VectorStoreRetriever
from langchain.chains import RetrievalQA
from langchain_community.chat_models.sambanova import ChatSambaNovaCloud
from langchain_core.prompts import PromptTemplate
from langchain_core.embeddings import Embeddings
import faiss

API_KEY = os.getenv("SAMBA_API_KEY", "540f8914-997e-46c6-829a-ff76f5d4d265")
API_URL = os.getenv("SAMBA_API_URL", "https://api.sambanova.ai/v1/chat/completions")

embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

class SentenceTransformerEmbeddings(Embeddings):
    """Wrapper around SentenceTransformer for compatibility with LangChain."""
    def __init__(self, model):
        self.model = model

    def embed_documents(self, texts):
        return self.model.encode(texts, show_progress_bar=True, convert_to_tensor=True)

    def embed_query(self, text):
        return self.model.encode(text, convert_to_tensor=True)

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file."""
    text = ""
    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting PPT: {e}")
    return text

def chunk_text(text, max_chunk_size=1000, overlap=200):
    """Split text into chunks with overlap."""
    chunks = []
    for i in range(0, len(text), max_chunk_size - overlap):
        chunks.append(text[i:i + max_chunk_size])
    return chunks

def create_faiss_index(content_chunks):
    """Create a FAISS index from content chunks."""
    embeddings = SentenceTransformerEmbeddings(embedding_model)
    return FAISS.from_texts(
        texts=content_chunks,
        embedding=embeddings,
        metadatas=[{"source": "ppt", "index": i} for i in range(len(content_chunks))]
    )

def save_content_and_vectors(content_chunks, faiss_index):
    """Save content and the FAISS index."""
    os.makedirs('extracted_content', exist_ok=True)
    with open('extracted_content/extracted_content.txt', 'w', encoding='utf-8') as f:
        for chunk in content_chunks:
            f.write(chunk + "\n\n")
    print("Extracted content saved.")

    os.makedirs('faiss_index', exist_ok=True)
    faiss.write_index(faiss_index.index, 'faiss_index/index.faiss')
    print("FAISS index saved.")

def get_custom_prompt():
    """Define a custom prompt template."""
    template = """<s>[INST] <<SYS>>
    You are a highly intelligent assistant. Use the provided context to answer the question.
    Do NOT use any information outside the given context. If the answer is not found in the context,
    respond with: "I couldn't find this information in the provided context."

    Context:
    {context}
    <</SYS>>

    Question: {question}
    Answer strictly based on the context provided.
    [/INST]"""
    return PromptTemplate(template=template, input_variables=["context", "question"])

def main():
    ppt_path = input("Enter the path to the PPT file: ").strip()

    ppt_text = extract_text_from_ppt(ppt_path)
    if not ppt_text.strip():
        print("No text extracted.")
        return

    print("PPT Text Extracted Successfully.")
    content_chunks = chunk_text(ppt_text, max_chunk_size=1000, overlap=200)
    print(f"Extracted {len(content_chunks)} content chunks.")

    faiss_index = create_faiss_index(content_chunks)
    print("FAISS index created.")

    save_content_and_vectors(content_chunks, faiss_index)

    retriever = VectorStoreRetriever(vectorstore=faiss_index, search_kwargs={"k": 5})

    llm = ChatSambaNovaCloud(
        model="llama3-405b",
        max_tokens=2048,
        temperature=0.5,
        top_k=10,
        top_p=0.9,
        request_timeout=30
    )

    custom_prompt = get_custom_prompt()
    qa_chain = RetrievalQA.from_chain_type(
        llm=llm,
        retriever=retriever,
        chain_type_kwargs={"prompt": custom_prompt, "document_variable_name": "context"}
    )

    while True:
        query = input("Ask a question (or type 'exit'): ").strip()
        if query.lower() == 'exit':
            print("Thank U!")
            break

        try:
            retrieved_docs = retriever.get_relevant_documents(query)
            context = "\n".join([doc.page_content for doc in retrieved_docs])

            if not context.strip():
                print("I couldn't find this information in the provided context.")
                continue

            result = qa_chain.invoke({"query": query, "context": context})
            print(f"Response: {result['result']}")
        except Exception as e:
            print(f"Error: {e}")

if __name__ == "__main__":
    main()
