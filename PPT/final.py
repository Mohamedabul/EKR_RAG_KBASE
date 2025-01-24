import os
import requests
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from langchain_community.vectorstores import FAISS
from langchain_core.vectorstores import VectorStoreRetriever
from langchain.chains import RetrievalQA
from langchain_community.chat_models.sambanova import ChatSambaNovaCloud
from langchain_core.prompts import PromptTemplate
from langchain_core.embeddings import Embeddings
import faiss
import pymongo
from pymongo import MongoClient
from langchain.memory import ConversationBufferMemory

MONGO_URI = "mongodb+srv://2217028:NquBh8oPPopA0Zuu@sumrag.ux9hs.mongodb.net/?retryWrites=true&w=majority&appName=SUMRAG"
DB_NAME = "DT3_EKR_BASE"
COLLECTION_NAME = "vectors"

client = MongoClient(MONGO_URI)
db = client[DB_NAME]
collection = db[COLLECTION_NAME]

API_KEY = os.getenv("SAMBA_API_KEY", "540f8914-997e-46c6-829a-ff76f5d4d265")
API_URL = os.getenv("SAMBA_API_URL", "https://api.sambanova.ai/v1/chat/completions")

embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

class SentenceTransformerEmbeddings(Embeddings):
    def __init__(self, model):
        self.model = model

    def embed_documents(self, texts):
        return self.model.encode(texts, show_progress_bar=True, convert_to_tensor=True)

    def embed_query(self, text):
        return self.model.encode(text, convert_to_tensor=True)

def extract_text_from_ppt(ppt_path):
    text = ""
    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting PPT: {e}")
    return text

def summarize_text(text):
    """Uses SambaNova API to summarize the given text."""
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": "llama3-405b",
        "messages": [{"role": "user", "content": text}],
        "max_tokens": 1024,
        "temperature": 0.7,
        "top_k": 1,
        "top_p": 0.01
    }

    response = requests.post(API_URL, json=payload, headers=headers, timeout=30)

    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    else:
        raise Exception(f"Error: {response.status_code} - {response.text}")

def chunk_text(text, max_chunk_size=2000, overlap=1500):
    chunks = []
    for i in range(0, len(text), max_chunk_size - overlap):
        chunks.append(text[i:i + max_chunk_size])
    return chunks

def create_faiss_index(content_chunks):
    embeddings = SentenceTransformerEmbeddings(embedding_model)
    return FAISS.from_texts(
        texts=content_chunks,
        embedding=embeddings,
        metadatas=[{"source": "ppt", "index": i} for i in range(len(content_chunks))]
    )

def get_custom_prompt():
    custom_prompt_template = """<s>[INST] <<SYS>>
    You are an intelligent assistant designed to extract and provide comprehensive, detailed answers based on the content of a document.
    Use the following context, extracted from the file, to answer the user's questions. When responding, locate and include all relevant information related to the question from the document.

    Context:
    {context}

    <</SYS>>

    Please answer the following question:
    {question}

    Ensure your response includes all available and relevant content from the document to fully address the user's query. If the exact answer is not found within the data, respond with:
    'This information isn't available in the provided data.'
    [/INST]"""
    return PromptTemplate(template=custom_prompt_template, input_variables=['context', 'question'])

def main():
    ppt_path = input("Enter the path to the PPT file: ").strip()
    ppt_text = extract_text_from_ppt(ppt_path)
    if not ppt_text.strip():
        print("No text extracted.")
        return

    print("PPT Text Extracted Successfully.")

    summary = summarize_text(ppt_text)
    print("Summary:\n", summary)

    content_chunks = chunk_text(ppt_text, max_chunk_size=1000, overlap=200)
    print(f"Extracted {len(content_chunks)} content chunks.")

    faiss_index = create_faiss_index(content_chunks)
    print("FAISS index created.")

    retriever = VectorStoreRetriever(vectorstore=faiss_index, search_kwargs={"k": 5})
    llm = ChatSambaNovaCloud(
        model="llama3-405b",
        max_tokens=2048,
        temperature=0.5,
        top_k=10,
        top_p=0.9,
        request_timeout=30
    )
    custom_prompt = get_custom_prompt()
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
    qa_chain = RetrievalQA.from_chain_type(
        llm=llm,
        retriever=retriever,
        chain_type_kwargs={"prompt": custom_prompt, "document_variable_name": "context"},
        memory=memory
    )
    user_interactions = []
    while True:
        query = input("Ask a question (or type 'exit'): ").strip()
        if query.lower() == 'exit':
            print("Thank you! Exiting and saving data.")
            break

        try:
            retrieved_docs = retriever.invoke(query)
            context = "\n".join([doc.page_content for doc in retrieved_docs])

            if not context.strip():
                print("I couldn't find this information in the provided context.")
                continue

            result = qa_chain.invoke({"query": query, "context": context})
            response = result['result']
            print(f"Response: {response}")
            user_interactions.append({
                "query": query,
                "response": response
            })
        except Exception as e:
            print(f"Error: {e}")

    data_to_store = {
        "summary": summary,
        "interactions": user_interactions
    }
    collection.insert_one(data_to_store)
    print("Data saved to MongoDB.")

if __name__ == "__main__":
    main()
